"""Read-only structural benchmark for local document samples.

The script intentionally writes no files and never prints input paths or content.
It is documentation experiment tooling, not an ArcheOS Adapter or runtime.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import subprocess
from time import perf_counter
from zipfile import ZipFile


def _pdf(path: Path) -> dict[str, int]:
    import pdfplumber

    with pdfplumber.open(path) as document:
        return {
            "pages": len(document.pages),
            "words": sum(len(page.extract_words()) for page in document.pages),
            "tables": sum(len(page.find_tables()) for page in document.pages),
        }


def _xlsx(path: Path) -> dict[str, int]:
    from openpyxl import load_workbook

    with ZipFile(path) as archive:
        members = archive.namelist()
    workbook = load_workbook(path, read_only=True, data_only=False)
    try:
        return {
            "sheets": len(workbook.worksheets),
            "hidden_sheets": sum(
                sheet.sheet_state != "visible" for sheet in workbook.worksheets
            ),
            "image_parts": sum(member.startswith("xl/media/") for member in members),
            "drawing_parts": sum(member.startswith("xl/drawings/") for member in members),
        }
    finally:
        workbook.close()


def _shape_count(shapes: object) -> int:
    total = 0
    for shape in shapes:  # type: ignore[union-attr]
        total += 1
        if hasattr(shape, "shapes"):
            total += _shape_count(shape.shapes)
    return total


def _pptx(path: Path) -> dict[str, int]:
    from pptx import Presentation

    presentation = Presentation(path)
    slides = presentation.slides
    return {
        "slides": len(slides),
        "shapes": sum(_shape_count(slide.shapes) for slide in slides),
        "notes_slides": sum(slide.has_notes_slide for slide in slides),
    }


def _image(path: Path) -> dict[str, str]:
    """Return only non-content structure fields from local system tools."""
    mime_type = subprocess.run(
        ["file", "--brief", "--mime-type", str(path)],
        check=True,
        capture_output=True,
        text=True,
    ).stdout.strip()
    sips = subprocess.run(
        ["sips", "-g", "format", "-g", "pixelWidth", "-g", "pixelHeight", "-g", "alphaInfo", str(path)],
        check=True,
        capture_output=True,
        text=True,
    ).stdout.splitlines()
    fields: dict[str, str] = {"mime_type": mime_type}
    for line in sips:
        for name in ("format", "pixelWidth", "pixelHeight", "alphaInfo"):
            marker = f"{name}: "
            if line.strip().startswith(marker):
                fields[name] = line.strip().split(marker, 1)[1]
    return fields


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--format", choices=("pdf", "xlsx", "pptx", "image"), required=True
    )
    parser.add_argument("input", type=Path)
    args = parser.parse_args()
    handlers = {"pdf": _pdf, "xlsx": _xlsx, "pptx": _pptx, "image": _image}
    started = perf_counter()
    result = handlers[args.format](args.input)
    result["elapsed_milliseconds"] = round((perf_counter() - started) * 1000)
    print(json.dumps(result, sort_keys=True))
    return 0


if __name__ == "__main__":  # pragma: no cover - command-line entry point
    raise SystemExit(main())
