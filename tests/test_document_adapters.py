from __future__ import annotations

import io
import json
import tempfile
import unittest
from contextlib import redirect_stderr, redirect_stdout
from pathlib import Path

from archeos.cli import main
from archeos.representation import LocalRepresentationRepository, RepresentationError, RepresentationService
from archeos.representation.adapters import (
    ImagePreflightRepresentationAdapter,
    MarkdownRepresentationAdapter,
    PdfTextRepresentationAdapter,
    PptxRepresentationAdapter,
    XlsxRepresentationAdapter,
)
from archeos.representation.registry import production_adapter
from archeos.source import LocalManagedSourceRepository


TIMESTAMP = "2026-08-13T00:00:00.000Z"


def _write_minimal_pdf(path: Path, text: str | None) -> None:
    stream = "" if text is None else f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET"
    objects = [
        "<< /Type /Catalog /Pages 2 0 R >>",
        "<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        "<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        f"<< /Length {len(stream.encode('ascii'))} >>\nstream\n{stream}\nendstream",
    ]
    result = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, payload in enumerate(objects, start=1):
        offsets.append(len(result))
        result.extend(f"{index} 0 obj\n{payload}\nendobj\n".encode("ascii"))
    xref = len(result)
    result.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    result.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        result.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    result.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode(
            "ascii"
        )
    )
    path.write_bytes(result)


class DocumentAdapterTest(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)
        self.managed_root = self.root / "managed"
        self.representation_root = self.root / "representations"
        self.source_number = 0

    def tearDown(self) -> None:
        self.temp.cleanup()

    def build(
        self,
        external: Path,
        media_type: str,
        adapter: object,
        configuration: dict[str, object] | None = None,
    ):
        self.source_number += 1
        source_id = f"src_{self.source_number:032x}"
        source_repository = LocalManagedSourceRepository(
            self.managed_root,
            id_factory=lambda: source_id,
            clock=lambda: TIMESTAMP,
        )
        source = source_repository.admit(external, metadata={"media_type": media_type}).source
        service = RepresentationService(
            source_repository,
            LocalRepresentationRepository(self.representation_root, clock=lambda: TIMESTAMP),
            clock=lambda: TIMESTAMP,
        )
        result = service.build(source.source_id, adapter, configuration or {})
        artifact = result.representation.artifacts[0]
        payload = json.loads(
            (self.representation_root / source.source_id / result.representation.representation_id / artifact.locator).read_text(
                encoding="utf-8"
            )
        )
        return source, service, result, payload

    def test_markdown_preserves_authored_structure_and_line_locators(self) -> None:
        external = self.root / "fixture.md"
        external.write_text(
            "---\ntitle: synthetic\n---\n# Title\n\nParagraph [link](https://invalid.test) [[local]]\n\n> Quote\n\n- item\n\n```txt\ncode\n```\n\n| A | B |\n| - | - |\n| 1 | 2 |\n",
            encoding="utf-8",
        )
        _, _, result, payload = self.build(external, "text/markdown", MarkdownRepresentationAdapter())
        self.assertEqual(result.representation.status, "complete")
        kinds = {block["kind"] for block in payload["blocks"]}
        self.assertTrue({"frontmatter", "heading", "paragraph", "quote", "list_item", "code_block", "table"} <= kinds)
        paragraph = next(item for item in payload["blocks"] if item["kind"] == "paragraph")
        self.assertEqual(paragraph["source_locator"], {"line_start": 6, "line_end": 6})
        self.assertEqual(paragraph["authored_links"], ["[link](https://invalid.test)", "[[local]]"])

    def test_pdf_text_has_page_bbox_locators_and_blank_pdf_is_partial(self) -> None:
        text_pdf = self.root / "text.pdf"
        _write_minimal_pdf(text_pdf, "Synthetic PDF")
        _, _, result, payload = self.build(text_pdf, "application/pdf", PdfTextRepresentationAdapter())
        self.assertEqual(result.representation.status, "partial")
        block = payload["pages"][0]["text_blocks"][0]
        self.assertEqual(block["source_locator"]["page"], 1)
        self.assertEqual(len(block["source_locator"]["bbox"]), 4)
        self.assertIn(
            "READING_ORDER_NOT_VERIFIED", {item.code for item in result.representation.warnings}
        )

        blank_pdf = self.root / "blank.pdf"
        _write_minimal_pdf(blank_pdf, None)
        _, _, blank_result, _ = self.build(blank_pdf, "application/pdf", PdfTextRepresentationAdapter())
        self.assertEqual(blank_result.representation.status, "partial")
        self.assertIn("SCANNED_CONTENT_UNSUPPORTED", {item.code for item in blank_result.representation.warnings})

    def test_xlsx_preserves_formula_structure_and_visibility(self) -> None:
        from PIL import Image
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image as XlsxImage
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.worksheet.table import Table, TableStyleInfo

        external = self.root / "fixture.xlsx"
        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Synthetic"
        worksheet.append(["name", "value"])
        worksheet.append(["row", 2])
        worksheet["C2"] = "=B2+1"
        worksheet.merge_cells("D1:E1")
        worksheet.row_dimensions[2].hidden = True
        worksheet.column_dimensions["E"].hidden = True
        table = Table(displayName="SyntheticTable", ref="A1:B2")
        table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
        worksheet.add_table(table)
        formatted_blank = worksheet["G4"]
        formatted_blank.number_format = "0.00%"
        formatted_blank.font = Font(name="Synthetic Font", bold=True)
        formatted_blank.fill = PatternFill("solid", fgColor="FF00FF00")
        formatted_blank.alignment = Alignment(horizontal="center", wrap_text=True)
        image_path = self.root / "xlsx-image.png"
        Image.new("RGBA", (1, 1), (0, 0, 0, 0)).save(image_path)
        worksheet.add_image(XlsxImage(str(image_path)), "F3")
        workbook.create_sheet("Hidden").sheet_state = "hidden"
        workbook.save(external)
        _, _, result, payload = self.build(
            external,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            XlsxRepresentationAdapter(),
        )
        self.assertEqual(result.representation.status, "partial")
        sheet = payload["sheets"][0]
        formula = next(cell for cell in sheet["cells"] if cell["source_locator"]["cell"] == "C2")
        self.assertEqual(formula["formula"], "=B2+1")
        self.assertIsNone(formula["cached_value"])
        self.assertIn(
            "FORMULA_CACHE_UNAVAILABLE", {item.code for item in result.representation.warnings}
        )
        self.assertIn("D1:E1", sheet["merged_ranges"])
        self.assertIn(2, sheet["hidden_rows"])
        self.assertIn("E", sheet["hidden_columns"])
        self.assertEqual(sheet["tables"], [{"name": "SyntheticTable", "ref": "A1:B2"}])
        self.assertEqual(sheet["embedded_media"][0]["anchor"]["from"], {"row": 3, "column": 6})
        blank = next(cell for cell in sheet["cells"] if cell["source_locator"]["cell"] == "G4")
        self.assertIsNone(blank["value"])
        self.assertEqual(blank["blank_structure"]["number_format"], "0.00%")
        self.assertTrue(blank["blank_structure"]["font"]["bold"])
        self.assertEqual(blank["blank_structure"]["fill"]["foreground"], "FF00FF00")
        self.assertTrue(blank["blank_structure"]["alignment"]["wrap_text"])
        self.assertTrue(payload["sheets"][1]["hidden"])

    def test_pptx_preserves_slide_shape_text_and_table_locators(self) -> None:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        external = self.root / "fixture.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        text_box.text = "Synthetic slide"
        table = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(2), Inches(1)).table
        table.cell(0, 0).text = "table value"
        group = slide.shapes.add_group_shape()
        grouped_text = group.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        grouped_text.text = "Grouped synthetic slide"
        image_path = self.root / "pptx-image.png"
        Image.new("RGB", (1, 1), (1, 2, 3)).save(image_path)
        slide.shapes.add_picture(str(image_path), Inches(4), Inches(1), Inches(1), Inches(1))
        presentation.save(external)
        _, _, result, payload = self.build(
            external,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            PptxRepresentationAdapter(),
        )
        self.assertEqual(result.representation.status, "partial")
        shapes = payload["slides"][0]["shapes"]
        self.assertTrue(any(item.get("text") == "Synthetic slide" for item in shapes))
        table_shape = next(item for item in shapes if "table" in item)
        self.assertEqual(table_shape["table"], [["table value"]])
        self.assertIn("slide_id", table_shape["source_locator"])
        self.assertIn("shape_id", table_shape["source_locator"])
        self.assertTrue(any("media" in item for item in shapes))
        group_shape = next(item for item in shapes if "children" in item)
        child = group_shape["children"][0]
        self.assertEqual(child["text"], "Grouped synthetic slide")
        self.assertEqual(
            child["source_locator"]["group_path"], [group_shape["source_locator"]["shape_id"]]
        )
        self.assertIn(
            "GROUP_LAYOUT_COORDINATES_UNVERIFIED",
            {item.code for item in result.representation.warnings},
        )

    def test_image_preflight_is_local_and_privacy_route_is_explicit(self) -> None:
        external = self.root / "fixture.png"
        external.write_bytes(
            b"\x89PNG\r\n\x1a\n" + b"\x00\x00\x00\rIHDR" + (3).to_bytes(4, "big")
            + (2).to_bytes(4, "big") + b"\x08\x06\x00\x00\x00"
        )
        _, _, result, payload = self.build(
            external, "image/png", ImagePreflightRepresentationAdapter(), {"privacy_route": "restricted"}
        )
        self.assertEqual(result.representation.status, "complete")
        self.assertEqual(payload["privacy_route"], "restricted")
        self.assertEqual((payload["pixel_width"], payload["pixel_height"]), (3, 2))
        self.assertIn("RESTRICTED_LOCAL_ONLY", {item.code for item in result.representation.warnings})

    def test_unknown_image_route_is_partial_and_invalid_route_fails_closed(self) -> None:
        external = self.root / "fixture-unknown.png"
        external.write_bytes(b"not an image")
        _, service, result, _ = self.build(external, "image/png", ImagePreflightRepresentationAdapter())
        self.assertEqual(result.representation.status, "partial")
        self.assertIn("PRIVACY_ROUTE_UNKNOWN", {item.code for item in result.representation.warnings})
        with self.assertRaises(RepresentationError):
            service.build(result.representation.source_id, ImagePreflightRepresentationAdapter(), {"privacy_route": "cloud"})

    def test_adapter_runtime_error_publishes_nothing(self) -> None:
        external = self.root / "corrupt.pdf"
        external.write_bytes(b"not a PDF")
        self.source_number += 1
        source_id = f"src_{self.source_number:032x}"
        source_repository = LocalManagedSourceRepository(
            self.managed_root, id_factory=lambda: source_id, clock=lambda: TIMESTAMP
        )
        source = source_repository.admit(external, metadata={"media_type": "application/pdf"}).source
        service = RepresentationService(
            source_repository,
            LocalRepresentationRepository(self.representation_root, clock=lambda: TIMESTAMP),
            clock=lambda: TIMESTAMP,
        )
        with self.assertRaises(RepresentationError):
            service.build(source.source_id, PdfTextRepresentationAdapter(), {})
        self.assertFalse((self.representation_root / source.source_id).exists())
        staging = self.representation_root / ".staging"
        self.assertFalse(staging.exists() and any(staging.iterdir()))

    def test_registry_and_cli_register_only_approved_adapters(self) -> None:
        self.assertEqual(production_adapter("markdown").version, "4.2.0")
        self.assertEqual(production_adapter("pdf-text").version, "0.11.10")
        self.assertEqual(production_adapter("xlsx").version, "3.1.5")
        self.assertEqual(production_adapter("pptx").version, "1.0.2")
        with self.assertRaises(Exception):
            production_adapter("ocr")

        external = self.root / "cli.md"
        external.write_text("# CLI synthetic\n", encoding="utf-8")
        managed_root = self.root / "cli-managed"
        representation_root = self.root / "cli-representations"
        output = io.StringIO()
        with redirect_stdout(output):
            self.assertEqual(
                main([
                    "source", "admit", str(external), "--managed-root", str(managed_root),
                    "--media-type", "text/markdown", "--source-id", "src_" + "a" * 32,
                ]),
                0,
            )
            self.assertEqual(
                main([
                    "representation", "build", "src_" + "a" * 32, "--adapter", "markdown",
                    "--managed-root", str(managed_root), "--representation-root", str(representation_root),
                ]),
                0,
            )
        self.assertIn('"status": "built"', output.getvalue())

    def test_cli_image_privacy_route_is_fingerprinted_and_fails_closed(self) -> None:
        external = self.root / "cli-image.png"
        external.write_bytes(
            b"\x89PNG\r\n\x1a\n" + b"\x00\x00\x00\rIHDR" + (3).to_bytes(4, "big")
            + (2).to_bytes(4, "big") + b"\x08\x06\x00\x00\x00"
        )
        managed_root = self.root / "cli-image-managed"
        representation_root = self.root / "cli-image-representations"
        source_id = "src_" + "b" * 32
        with redirect_stdout(io.StringIO()):
            self.assertEqual(
                main([
                    "source", "admit", str(external), "--managed-root", str(managed_root),
                    "--media-type", "image/png", "--source-id", source_id,
                ]),
                0,
            )

        def build(*arguments: str) -> tuple[int, dict[str, object]]:
            output = io.StringIO()
            with redirect_stdout(output):
                code = main([
                    "representation", "build", source_id, "--adapter", "image-preflight",
                    "--managed-root", str(managed_root), "--representation-root", str(representation_root),
                    *arguments,
                ])
            return code, json.loads(output.getvalue()) if code == 0 else {"error": output.getvalue()}

        default_code, default = build()
        unknown_code, unknown = build("--privacy-route", "unknown")
        restricted_code, restricted = build("--privacy-route", "restricted")
        self.assertEqual((default_code, unknown_code, restricted_code), (0, 0, 0))
        default_representation = default["representation"]["representation"]
        unknown_representation = unknown["representation"]["representation"]
        restricted_representation = restricted["representation"]["representation"]
        self.assertEqual(default_representation["status"], "partial")
        self.assertEqual(unknown["status"], "existing")
        self.assertEqual(
            default_representation["configuration_fingerprint"],
            unknown_representation["configuration_fingerprint"],
        )
        self.assertNotEqual(
            default_representation["configuration_fingerprint"],
            restricted_representation["configuration_fingerprint"],
        )
        self.assertEqual(restricted_representation["status"], "complete")
        with redirect_stderr(io.StringIO()), self.assertRaises(SystemExit):
            main([
                "representation", "build", source_id, "--adapter", "image-preflight",
                "--managed-root", str(managed_root), "--representation-root", str(representation_root),
                "--privacy-route", "cloud",
            ])
        manifests = list((representation_root / source_id).glob("*/manifest.json"))
        self.assertEqual(len(manifests), 2)
