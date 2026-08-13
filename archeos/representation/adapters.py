"""Approved local format Adapters isolated behind RepresentationAdapter."""

from __future__ import annotations

import json
import re
import warnings as standard_warnings
from collections.abc import Mapping
from pathlib import Path
from typing import Any

from ..source.models import ManagedSource
from .models import AdapterArtifact, AdapterBuildResult, RepresentationWarning


def _write_artifact(staging_dir: Path, name: str, payload: object) -> AdapterArtifact:
    locator = f"artifacts/{name}"
    path = staging_dir / locator
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, sort_keys=True) + "\n", encoding="utf-8")
    return AdapterArtifact("structure", locator, "application/json")


def _warning(code: str, message: str, severity: str = "warning") -> RepresentationWarning:
    return RepresentationWarning(code, message, severity)


def _json_value(value: object) -> object:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return str(value)


class MarkdownRepresentationAdapter:
    name = "markdown"
    version = "4.2.0"
    kind = "markdown_blocks"
    supported_media_types = ("text/markdown", "text/x-markdown", "text/plain")

    def build(
        self,
        source: ManagedSource,
        materialized_path: Path,
        staging_dir: Path,
        configuration: Mapping[str, object],
    ) -> AdapterBuildResult:
        from markdown_it import MarkdownIt

        text = materialized_path.read_text(encoding="utf-8")
        lines = text.splitlines(keepends=True)
        parser = MarkdownIt("commonmark").enable("table")
        blocks: list[dict[str, object]] = []
        start = 0
        if text.startswith("---\n") or text.startswith("---\r\n"):
            for index, line in enumerate(lines[1:], start=1):
                if line.rstrip("\r\n") in {"---", "..."}:
                    blocks.append(
                        {
                            "kind": "frontmatter",
                            "source_locator": {"line_start": 1, "line_end": index + 1},
                            "raw": "".join(lines[: index + 1]),
                        }
                    )
                    start = index + 1
                    break
        kind_map = {
            "heading_open": "heading",
            "paragraph_open": "paragraph",
            "list_item_open": "list_item",
            "blockquote_open": "quote",
            "fence": "code_block",
            "code_block": "code_block",
            "table_open": "table",
        }
        for token in parser.parse(text):
            kind = kind_map.get(token.type)
            if kind is None or token.map is None:
                continue
            line_start, line_end = token.map
            if line_end <= start:
                continue
            raw = "".join(lines[line_start:line_end])
            authored = re.findall(r"\[\[[^\]]+\]\]|\[[^\]]*\]\([^)]*\)", raw)
            blocks.append(
                {
                    "kind": kind,
                    "source_locator": {"line_start": line_start + 1, "line_end": line_end},
                    "raw": raw,
                    "authored_links": authored,
                }
            )
        artifact = _write_artifact(staging_dir, "blocks.json", {"blocks": blocks})
        return AdapterBuildResult(self.kind, (artifact,), 1.0)


class PdfTextRepresentationAdapter:
    name = "pdf-text"
    version = "0.11.10"
    kind = "pdf_text"
    supported_media_types = ("application/pdf",)

    def build(
        self,
        source: ManagedSource,
        materialized_path: Path,
        staging_dir: Path,
        configuration: Mapping[str, object],
    ) -> AdapterBuildResult:
        import pdfplumber

        pages: list[dict[str, object]] = []
        warnings: list[RepresentationWarning] = []
        with pdfplumber.open(materialized_path) as document:
            for index, page in enumerate(document.pages, start=1):
                words = page.extract_words() or []
                blocks = [
                    {
                        "source_locator": {
                            "page": index,
                            "ordinal": ordinal,
                            "bbox": [word["x0"], word["top"], word["x1"], word["bottom"]],
                            "page_width": page.width,
                            "page_height": page.height,
                        },
                        "text": word["text"],
                    }
                    for ordinal, word in enumerate(words, start=1)
                ]
                tables = []
                for ordinal, table in enumerate(page.find_tables() or [], start=1):
                    tables.append(
                        {
                            "source_locator": {
                                "page": index,
                                "ordinal": ordinal,
                                "bbox": list(table.bbox),
                                "page_width": page.width,
                                "page_height": page.height,
                            },
                            "cells": table.extract(),
                        }
                    )
                if not words:
                    warnings.append(
                        _warning(
                            "SCANNED_CONTENT_UNSUPPORTED",
                            f"page {index} has no extractable text; OCR is intentionally deferred",
                        )
                    )
                if page.images:
                    warnings.append(
                        _warning(
                            "IMAGE_CONTENT_NOT_EXTRACTED",
                            f"page {index} contains image content not represented by the text Adapter",
                        )
                    )
                if words:
                    warnings.append(
                        _warning(
                            "READING_ORDER_NOT_VERIFIED",
                            f"page {index} word order is extraction order, not a verified reading order",
                        )
                    )
                pages.append(
                    {
                        "source_locator": {
                            "page": index,
                            "page_width": page.width,
                            "page_height": page.height,
                        },
                        "text_blocks": blocks,
                        "tables": tables,
                    }
                )
        completeness = 1.0 if not warnings else 0.5
        artifact = _write_artifact(staging_dir, "pages.json", {"pages": pages})
        return AdapterBuildResult(self.kind, (artifact,), completeness, tuple(warnings))


class XlsxRepresentationAdapter:
    name = "xlsx"
    version = "3.1.5"
    kind = "xlsx_structure"
    supported_media_types = (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

    def build(
        self,
        source: ManagedSource,
        materialized_path: Path,
        staging_dir: Path,
        configuration: Mapping[str, object],
    ) -> AdapterBuildResult:
        from openpyxl import load_workbook

        max_cells = configuration.get("max_cells", 100_000)
        if isinstance(max_cells, bool) or not isinstance(max_cells, int) or max_cells <= 0:
            raise ValueError("max_cells must be a positive integer")
        warnings: list[RepresentationWarning] = []
        with standard_warnings.catch_warnings(record=True) as load_warnings:
            standard_warnings.simplefilter("always")
            workbook = load_workbook(
                materialized_path, read_only=False, data_only=False, keep_links=False
            )
            cached = load_workbook(
                materialized_path, read_only=False, data_only=True, keep_links=False
            )
        if load_warnings:
            warnings.append(
                _warning(
                    "EMBEDDED_MEDIA_UNSUPPORTED",
                    "workbook reader reported unsupported embedded content; it was not retained",
                )
            )
        sheets: list[dict[str, object]] = []
        remaining = max_cells
        try:
            for index, worksheet in enumerate(workbook.worksheets, start=1):
                cached_sheet = cached[worksheet.title]
                cells: list[dict[str, object]] = []
                for row in worksheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        if remaining == 0:
                            warnings.append(
                                _warning("CELL_LIMIT_REACHED", "configured cell limit reached")
                            )
                            break
                        formula = cell.value if cell.data_type == "f" else None
                        cached_value = (
                            _json_value(cached_sheet[cell.coordinate].value)
                            if formula is not None
                            else None
                        )
                        if formula is not None and cached_value is None:
                            warnings.append(
                                _warning(
                                    "FORMULA_CACHE_UNAVAILABLE",
                                    f"{worksheet.title}!{cell.coordinate} has no cached value",
                                )
                            )
                        cells.append(
                            {
                                "source_locator": {
                                    "sheet": worksheet.title,
                                    "sheet_index": index,
                                    "cell": cell.coordinate,
                                },
                                "value": _json_value(cell.value),
                                "formula": formula,
                                "cached_value": cached_value,
                            }
                        )
                        remaining -= 1
                    if remaining == 0:
                        break
                media = []
                for media_index, image in enumerate(getattr(worksheet, "_images", ()), start=1):
                    anchor = self._image_anchor(image)
                    if anchor is None:
                        warnings.append(
                            _warning(
                                "MEDIA_ANCHOR_UNAVAILABLE",
                                f"{worksheet.title} media {media_index} has no stable cell anchor",
                            )
                        )
                    media.append(
                        {
                            "source_locator": {
                                "sheet": worksheet.title,
                                "sheet_index": index,
                                "media_index": media_index,
                            },
                            "anchor": anchor,
                        }
                    )
                sheets.append(
                    {
                        "source_locator": {"sheet": worksheet.title, "sheet_index": index},
                        "hidden": worksheet.sheet_state != "visible",
                        "cells": cells,
                        "merged_ranges": [str(item) for item in worksheet.merged_cells.ranges],
                        "hidden_rows": [key for key, item in worksheet.row_dimensions.items() if item.hidden],
                        "hidden_columns": [key for key, item in worksheet.column_dimensions.items() if item.hidden],
                        "tables": [
                            {"name": table.name, "ref": table.ref}
                            for table in worksheet.tables.values()
                        ],
                        "embedded_media": media,
                    }
                )
        finally:
            workbook.close()
            cached.close()
        completeness = 1.0 if not warnings else 0.8
        artifact = _write_artifact(staging_dir, "workbook.json", {"sheets": sheets})
        return AdapterBuildResult(self.kind, (artifact,), completeness, tuple(warnings))

    @staticmethod
    def _image_anchor(image: object) -> dict[str, object] | None:
        anchor = getattr(image, "anchor", None)
        start = getattr(anchor, "_from", None)
        if start is None:
            return None
        payload: dict[str, object] = {
            "from": {"row": start.row + 1, "column": start.col + 1},
        }
        end = getattr(anchor, "to", None)
        if end is not None:
            payload["to"] = {"row": end.row + 1, "column": end.col + 1}
        return payload


class PptxRepresentationAdapter:
    name = "pptx"
    version = "1.0.2"
    kind = "pptx_structure"
    supported_media_types = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    def build(
        self,
        source: ManagedSource,
        materialized_path: Path,
        staging_dir: Path,
        configuration: Mapping[str, object],
    ) -> AdapterBuildResult:
        from pptx import Presentation

        presentation = Presentation(materialized_path)
        warnings: list[RepresentationWarning] = []
        slides: list[dict[str, object]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            shapes: list[dict[str, object]] = []
            for shape_order, shape in enumerate(slide.shapes, start=1):
                locator = {
                    "slide_index": slide_index,
                    "slide_id": slide.slide_id,
                    "shape_id": shape.shape_id,
                    "shape_order": shape_order,
                    "shape_type": str(shape.shape_type),
                    "bbox": [shape.left, shape.top, shape.width, shape.height],
                }
                shape_data: dict[str, object] = {"source_locator": locator}
                if getattr(shape, "has_text_frame", False):
                    shape_data["text"] = shape.text
                if getattr(shape, "has_table", False):
                    shape_data["table"] = [
                        [cell.text for cell in row.cells] for row in shape.table.rows
                    ]
                if hasattr(shape, "image"):
                    image = shape.image
                    shape_data["media"] = {
                        "format": str(image.ext),
                        "content_type": image.content_type,
                        "size_bytes": len(image.blob),
                    }
                shapes.append(shape_data)
            notes = None
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
            xml = slide._element.xml
            if "<p:timing" in xml:
                warnings.append(
                    _warning("ANIMATION_UNSUPPORTED", f"slide {slide_index} has animation timing")
                )
            if "<p:oleObj" in xml:
                warnings.append(_warning("EMBEDDED_OBJECT_UNSUPPORTED", f"slide {slide_index} has embedded object"))
            slides.append(
                {
                    "source_locator": {"slide_index": slide_index, "slide_id": slide.slide_id},
                    "shapes": shapes,
                    "speaker_notes": notes,
                }
            )
        completeness = 1.0 if not warnings else 0.8
        artifact = _write_artifact(staging_dir, "slides.json", {"slides": slides})
        return AdapterBuildResult(self.kind, (artifact,), completeness, tuple(warnings))


class ImagePreflightRepresentationAdapter:
    name = "image-preflight"
    version = "1.0"
    kind = "image_structural_preflight"
    supported_media_types = ("image/jpeg", "image/png", "image/gif")

    def build(
        self,
        source: ManagedSource,
        materialized_path: Path,
        staging_dir: Path,
        configuration: Mapping[str, object],
    ) -> AdapterBuildResult:
        route = configuration.get("privacy_route", "unknown")
        if route not in {"unknown", "restricted", "standard"}:
            raise ValueError("privacy_route must be unknown, restricted, or standard")
        with materialized_path.open("rb") as handle:
            header = handle.read(32)
        structure: dict[str, Any]
        warnings: list[RepresentationWarning] = []
        if header.startswith(b"\x89PNG\r\n\x1a\n") and len(header) >= 26:
            width = int.from_bytes(header[16:20], "big")
            height = int.from_bytes(header[20:24], "big")
            color_type = header[25]
            structure = {
                "format": "png",
                "pixel_width": width,
                "pixel_height": height,
                "alpha": color_type in {4, 6},
            }
        elif header.startswith(b"GIF87a") or header.startswith(b"GIF89a"):
            structure = {
                "format": "gif",
                "pixel_width": int.from_bytes(header[6:8], "little"),
                "pixel_height": int.from_bytes(header[8:10], "little"),
                "alpha": "unknown",
            }
            warnings.append(
                _warning("ALPHA_NOT_AVAILABLE", "GIF alpha cannot be determined by local preflight")
            )
        elif header.startswith(b"\xff\xd8"):
            width, height = self._jpeg_dimensions(materialized_path)
            structure = {"format": "jpeg", "pixel_width": width, "pixel_height": height, "alpha": False}
        else:
            structure = {"format": "unknown", "pixel_width": None, "pixel_height": None, "alpha": "unknown"}
            warnings.append(
                _warning(
                    "UNSUPPORTED_IMAGE_FORMAT",
                    "image format is not supported by local preflight",
                )
            )
        structure["media_type"] = source.media_type
        structure["privacy_route"] = route
        structure["orientation"] = "not_collected"
        if route == "unknown":
            warnings.append(_warning("PRIVACY_ROUTE_UNKNOWN", "privacy route was not explicitly selected"))
        if route == "restricted":
            warnings.append(
                _warning(
                    "RESTRICTED_LOCAL_ONLY",
                    "restricted input remains local; OCR and cloud routes are disabled",
                    "info",
                )
            )
        completeness = 1.0 if not any(item.severity == "warning" for item in warnings) else 0.8
        artifact = _write_artifact(staging_dir, "preflight.json", structure)
        return AdapterBuildResult(self.kind, (artifact,), completeness, tuple(warnings))

    @staticmethod
    def _jpeg_dimensions(path: Path) -> tuple[int | None, int | None]:
        with path.open("rb") as handle:
            if handle.read(2) != b"\xff\xd8":
                return None, None
            while True:
                marker_start = handle.read(1)
                if not marker_start:
                    break
                if marker_start != b"\xff":
                    continue
                marker = handle.read(1)
                while marker == b"\xff":
                    marker = handle.read(1)
                if marker in {b"\xd8", b"\xd9"}:
                    continue
                length_bytes = handle.read(2)
                if len(length_bytes) != 2:
                    break
                length = int.from_bytes(length_bytes, "big")
                if length < 2:
                    break
                if marker in {
                    b"\xc0", b"\xc1", b"\xc2", b"\xc3", b"\xc5", b"\xc6", b"\xc7",
                    b"\xc9", b"\xca", b"\xcb", b"\xcd", b"\xce", b"\xcf",
                }:
                    data = handle.read(length - 2)
                    if len(data) >= 5:
                        return int.from_bytes(data[3:5], "big"), int.from_bytes(data[1:3], "big")
                    break
                handle.seek(length - 2, 1)
        return None, None
